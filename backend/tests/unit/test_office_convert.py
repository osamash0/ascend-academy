"""Tests for office_convert (PPTX→PDF) and the .pptx upload-validation path."""
import io

import pytest
from pptx import Presentation

from backend.services import office_convert, upload_service


def _build_pptx(n_slides: int = 2) -> bytes:
    prs = Presentation()
    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"Slide {i + 1}"
        s.placeholders[1].text_frame.text = f"Body text {i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── office_convert ────────────────────────────────────────────────────────────

@pytest.mark.skipif(not office_convert.is_available(), reason="LibreOffice not installed")
async def test_to_pdf_produces_valid_pdf_with_matching_pages():
    pdf = await office_convert.to_pdf(_build_pptx(3), "deck.pptx")
    assert pdf[:4] == b"%PDF"
    import fitz
    with fitz.open(stream=pdf, filetype="pdf") as doc:
        assert len(doc) == 3  # slide N ↔ PDF page N


def test_find_soffice_honors_env_override(monkeypatch, tmp_path):
    fake = tmp_path / "soffice"
    fake.write_text("#!/bin/sh\n")
    monkeypatch.setenv("SOFFICE_BINARY", str(fake))
    assert office_convert._find_soffice() == str(fake)


def test_soffice_command_uses_launchservices_for_macos_app(monkeypatch):
    monkeypatch.setattr(office_convert.sys, "platform", "darwin")

    command = office_convert._soffice_command(
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        profile_dir="/tmp/profile", output_dir="/tmp/output", input_path="/tmp/deck.pptx",
    )

    assert command[:5] == ["open", "-W", "-a", "/Applications/LibreOffice.app", "--args"]
    assert "--headless" in command
    assert "/tmp/deck.pptx" in command


def test_soffice_command_stays_direct_off_macos(monkeypatch):
    monkeypatch.setattr(office_convert.sys, "platform", "linux")

    command = office_convert._soffice_command(
        "/usr/bin/soffice",
        profile_dir="/tmp/profile", output_dir="/tmp/output", input_path="/tmp/deck.pptx",
    )

    assert command[0] == "/usr/bin/soffice"
    assert command[1] == "--headless"


async def test_to_pdf_raises_without_soffice(monkeypatch):
    """Backstop only — callers reject .pptx via is_available() first.

    The message is deliberately the user-facing UNAVAILABLE_MESSAGE rather
    than the old "run brew install --cask libreoffice" text, which leaked
    developer instructions to students when it surfaced mid-upload.
    """
    monkeypatch.setattr(office_convert, "_find_soffice", lambda: None)
    with pytest.raises(RuntimeError, match="export your slides to PDF"):
        await office_convert.to_pdf(b"x", "deck.pptx")


# ── validate_upload (.pptx branch) ────────────────────────────────────────────

@pytest.fixture
def converter_available(monkeypatch):
    """Force the LibreOffice converter to look present.

    validate_upload now rejects .pptx outright when soffice is missing, so
    tests of the *container validation* logic must pin availability or they
    pass on a dev Mac with LibreOffice installed and fail in CI without it.
    """
    monkeypatch.setattr(office_convert, "is_available", lambda: True)


async def test_validate_upload_counts_slides(converter_available):
    assert await upload_service.validate_upload("lecture.pptx", _build_pptx(4)) == 4


async def test_validate_upload_rejects_bad_pptx(converter_available):
    with pytest.raises(ValueError):
        await upload_service.validate_upload("lecture.pptx", b"not a real pptx")


# ── the fail-fast guard ───────────────────────────────────────────────────────

async def test_validate_upload_rejects_pptx_when_converter_missing(monkeypatch):
    """The bug this closes: every Docker deployment ships curl only, so a
    .pptx upload was accepted, hashed, stored and enqueued, then failed later
    inside office_convert.to_pdf — mid-SSE, with a developer-facing message.
    """
    monkeypatch.setattr(office_convert, "is_available", lambda: False)

    with pytest.raises(ValueError, match="export your slides to PDF"):
        await upload_service.validate_upload("lecture.pptx", _build_pptx(2))


async def test_pdf_uploads_unaffected_when_converter_missing(monkeypatch, sample_pdf_bytes):
    """Control: the guard must be scoped to .pptx and never touch PDFs."""
    monkeypatch.setattr(office_convert, "is_available", lambda: False)

    assert await upload_service.validate_upload("lecture.pdf", sample_pdf_bytes) >= 1


def test_accepted_extensions_drops_pptx_when_converter_missing(monkeypatch):
    """/upload/config feeds the frontend's file picker. Advertising .pptx on a
    server that cannot convert it is what let the UI offer a file type whose
    every upload was guaranteed to fail."""
    from backend.api.v1 import upload as upload_api

    monkeypatch.setattr(office_convert, "is_available", lambda: False)
    assert upload_api.accepted_upload_extensions() == [".pdf"]

    monkeypatch.setattr(office_convert, "is_available", lambda: True)
    assert upload_api.accepted_upload_extensions() == [".pdf", ".pptx"]


async def test_validate_upload_rejects_unsupported_extension():
    with pytest.raises(ValueError, match="PDF and PowerPoint"):
        await upload_service.validate_upload("notes.txt", b"hello world")
