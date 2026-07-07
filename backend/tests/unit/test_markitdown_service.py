"""Unit tests for the MarkItDown parser service (PPTX → per-slide pages)."""
import io

import pytest
from pptx import Presentation
from pptx.util import Inches

from backend.services import markitdown_service as svc


def _build_pptx() -> bytes:
    """A 3-slide deck: title slide, bullet slide, and a blank slide with a
    bare textbox (no title placeholder)."""
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Introduction to Thermodynamics"
    s1.placeholders[1].text = "Lecture 1 — Prof. Mueller"

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "The First Law"
    body = s2.placeholders[1].text_frame
    body.text = "Energy is conserved"
    body.add_paragraph().text = "dU = δQ - δW"
    body.add_paragraph().text = "Internal energy is a state function"

    s3 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    tb = s3.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tb.text_frame.text = "Entropy always increases in an isolated system."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def test_extract_pages_splits_by_slide():
    pages = await svc.extract_pages(_build_pptx(), "lecture.pptx")

    # One entry per slide, 1-indexed.
    assert sorted(pages.keys()) == [1, 2, 3]

    # Titles come from the slide title placeholder (Markdown "# heading").
    assert pages[1]["title"] == "Introduction to Thermodynamics"
    assert pages[2]["title"] == "The First Law"

    # Body text is captured per slide.
    assert "Lecture 1" in pages[1]["text"]
    assert "Energy is conserved" in pages[2]["text"]
    assert "dU = δQ - δW" in pages[2]["text"]

    # A blank-layout slide with no title placeholder → title is None, text kept.
    assert pages[3]["title"] is None
    assert "Entropy always increases" in pages[3]["text"]


async def test_extract_pages_returns_expected_shape():
    pages = await svc.extract_pages(_build_pptx(), "lecture.pptx")
    for entry in pages.values():
        assert set(entry.keys()) == {"text", "title"}
        assert isinstance(entry["text"], str)
        assert entry["title"] is None or isinstance(entry["title"], str)


def test_parse_markdown_no_markers_falls_back_to_single_page():
    md = "# Some Document\n\nA paragraph of body text with no slide markers."
    pages = svc._parse_markdown_to_pages(md)
    assert list(pages.keys()) == [1]
    assert pages[1]["title"] == "Some Document"
    assert "body text" in pages[1]["text"]


def test_parse_markdown_empty_returns_nothing():
    assert svc._parse_markdown_to_pages("   \n  ") == {}


def test_title_skips_noise_lines():
    # First heading is a bullet/caption → no clean title.
    chunk = "# - just a dash bullet\nsome text"
    assert svc._title_from_chunk(chunk) is None
    # Too-short heading is ignored.
    assert svc._title_from_chunk("# ok\nbody") is None


async def test_extract_pages_raises_when_unavailable(monkeypatch):
    monkeypatch.setattr(svc, "_MARKITDOWN_AVAILABLE", False)
    with pytest.raises(RuntimeError, match="not installed"):
        await svc.extract_pages(b"x", "lecture.pptx")
