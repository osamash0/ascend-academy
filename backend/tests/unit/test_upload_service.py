"""Unit tests for backend.services.upload_service.

Covers the deterministic upload-pipeline plumbing that sits in front of the
parser: size capping, PDF/PPTX validation, storage upload (with bucket
auto-create), queue-depth backpressure, the Arq-enqueue + Redis pub/sub
streaming path, the production queue-outage refusal, the raw-page extractor's
per-parser branches, and the lazy import wrapper.

Everything external (Supabase storage, Arq/Redis, LibreOffice, alternate
parsers, the unified orchestrator) is mocked at its boundary — no network,
no real parse. Only fitz / python-pptx run for real, on tiny in-memory files.
"""
from __future__ import annotations

import io
import json

import pytest

from backend.services import upload_service
from backend.core import file_validation


# ── read_upload_capped ────────────────────────────────────────────────────────

class _FakeUpload:
    """Minimal UploadFile stand-in exposing async read(n)."""

    def __init__(self, data: bytes) -> None:
        self._data = data
        self._pos = 0

    async def read(self, n: int = -1) -> bytes:
        if n < 0:
            n = len(self._data) - self._pos
        chunk = self._data[self._pos : self._pos + n]
        self._pos += len(chunk)
        return chunk


async def test_read_upload_capped_returns_full_content_under_limit():
    data = b"x" * (3 * 1024 * 1024)  # 3 MB
    out = await upload_service.read_upload_capped(_FakeUpload(data), max_mb=5)
    assert out == data


async def test_read_upload_capped_raises_once_over_limit():
    data = b"y" * (2 * 1024 * 1024 + 10)  # just over 2 MB
    with pytest.raises(ValueError, match="exceeds the 2MB limit"):
        await upload_service.read_upload_capped(_FakeUpload(data), max_mb=2)


async def test_read_upload_capped_empty_file():
    out = await upload_service.read_upload_capped(_FakeUpload(b""), max_mb=1)
    assert out == b""


# ── validate_upload (PDF) ─────────────────────────────────────────────────────

async def test_validate_upload_pdf_returns_page_count(sample_pdf_bytes):
    assert await upload_service.validate_upload("lecture.pdf", sample_pdf_bytes) == 3


async def test_validate_upload_rejects_non_pdf_non_pptx():
    with pytest.raises(ValueError, match="Only PDF and PowerPoint"):
        await upload_service.validate_upload("notes.txt", b"%PDF-1.4 plus filler bytes")


async def test_validate_upload_corrupted_pdf(sample_pdf_bytes):
    # Passes magic-byte check (starts %PDF) but fitz can't open it → page_count -1.
    corrupt = b"%PDF-1.4 " + b"not a real pdf body " * 4
    with pytest.raises(ValueError, match="corrupted or password-protected"):
        await upload_service.validate_upload("broken.pdf", corrupt)


async def test_validate_upload_zero_page_pdf(monkeypatch, sample_pdf_bytes):
    import fitz

    class _FakeDoc:
        def __enter__(self):
            return self

        def __exit__(self, *a):
            return False

        def __len__(self):
            return 0

    monkeypatch.setattr(fitz, "open", lambda *a, **k: _FakeDoc())
    with pytest.raises(ValueError, match="PDF has no pages"):
        await upload_service.validate_upload("empty.pdf", sample_pdf_bytes)


async def test_validate_upload_too_many_pages(monkeypatch, sample_pdf_bytes):
    monkeypatch.setattr(upload_service, "MAX_PAGES", 2, raising=True)
    with pytest.raises(ValueError, match="Maximum 2 pages supported"):
        await upload_service.validate_upload("big.pdf", sample_pdf_bytes)


async def test_validate_upload_sanitizes_traversal_filename(sample_pdf_bytes):
    # A path-traversal filename is sanitized; extension check still resolves to PDF.
    assert await upload_service.validate_upload("../../etc/passwd.pdf", sample_pdf_bytes) == 3


# ── validate_upload / _validate_pptx (PPTX) ──────────────────────────────────

def _pptx_bytes(n_slides: int) -> bytes:
    from pptx import Presentation

    p = Presentation()
    for _ in range(n_slides):
        p.slides.add_slide(p.slide_layouts[6])  # blank layout
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


async def test_validate_upload_pptx_returns_slide_count():
    assert await upload_service.validate_upload("deck.pptx", _pptx_bytes(2)) == 2


async def test_validate_pptx_rejects_bad_magic_bytes():
    with pytest.raises(ValueError, match="Invalid PowerPoint file"):
        await upload_service.validate_upload("deck.pptx", b"NOTAZIP" + b"\x00" * 20)


async def test_validate_pptx_too_short():
    with pytest.raises(ValueError, match="Invalid PowerPoint file"):
        await upload_service.validate_upload("deck.pptx", b"PK")


async def test_validate_pptx_corrupted_zip():
    # Correct OOXML magic bytes but not a valid pptx → slide count -1.
    corrupt = b"PK\x03\x04" + b"garbage zip payload" * 4
    with pytest.raises(ValueError, match="corrupted or is not a valid"):
        await upload_service.validate_upload("deck.pptx", corrupt)


async def test_validate_pptx_zero_slides():
    with pytest.raises(ValueError, match="no slides"):
        await upload_service.validate_upload("deck.pptx", _pptx_bytes(0))


async def test_validate_pptx_too_many_slides(monkeypatch):
    monkeypatch.setattr(upload_service, "MAX_PAGES", 1, raising=True)
    with pytest.raises(ValueError, match="Maximum 1 slides supported"):
        await upload_service.validate_upload("deck.pptx", _pptx_bytes(2))


async def test_validate_pptx_exceeds_size_limit(monkeypatch):
    monkeypatch.setattr(file_validation, "MAX_FILE_BYTES", 100, raising=True)
    big = b"PK\x03\x04" + b"z" * 200
    with pytest.raises(ValueError, match="limit"):
        await upload_service.validate_upload("deck.pptx", big)


# ── queue_depth ──────────────────────────────────────────────────────────────

async def test_queue_depth_reads_zcard(monkeypatch):
    class _Pool:
        async def zcard(self, name):
            return 7

    async def _get_pool():
        return _Pool()

    monkeypatch.setattr(upload_service, "get_arq_pool", _get_pool)
    assert await upload_service.queue_depth() == 7


async def test_queue_depth_returns_zero_on_error(monkeypatch):
    async def _boom():
        raise RuntimeError("redis down")

    monkeypatch.setattr(upload_service, "get_arq_pool", _boom)
    assert await upload_service.queue_depth() == 0


# ── upload_pdf_to_storage ────────────────────────────────────────────────────

class _FakeStorageBucket:
    def __init__(self, fail_first: bool = False):
        self.uploads: list[tuple] = []
        self._fail_first = fail_first

    def upload(self, path, content, file_options=None):
        if self._fail_first and not self.uploads:
            # record nothing; simulate missing bucket on the first attempt
            self.uploads.append(("attempted-missing", path))
            raise RuntimeError("Bucket not found")
        self.uploads.append((path, len(content)))


class _FakeStorage:
    def __init__(self, bucket):
        self._bucket = bucket
        self.created_buckets: list[str] = []

    def from_(self, name):
        return self._bucket

    def create_bucket(self, name, options=None):
        self.created_buckets.append(name)


class _FakeSB:
    def __init__(self, storage):
        self.storage = storage


async def test_upload_pdf_to_storage_happy_path(monkeypatch):
    bucket = _FakeStorageBucket()
    storage = _FakeStorage(bucket)
    monkeypatch.setattr(upload_service, "get_client", lambda use_admin=False: _FakeSB(storage))

    from backend.core import database

    async def _run_sync(fn, *a, **k):
        return fn(*a, **k)

    monkeypatch.setattr(database, "run_sync", _run_sync)

    await upload_service.upload_pdf_to_storage("a" * 64, b"%PDF-1.4 data")
    assert bucket.uploads[-1][0] == "a" * 64 + ".pdf"


async def test_upload_pdf_to_storage_creates_missing_bucket(monkeypatch):
    bucket = _FakeStorageBucket(fail_first=True)
    storage = _FakeStorage(bucket)
    monkeypatch.setattr(upload_service, "get_client", lambda use_admin=False: _FakeSB(storage))

    from backend.core import database

    async def _run_sync(fn, *a, **k):
        return fn(*a, **k)

    monkeypatch.setattr(database, "run_sync", _run_sync)

    await upload_service.upload_pdf_to_storage("b" * 64, b"%PDF data")
    # It hit "Bucket not found", created the bucket, then retried the upload.
    assert "pdf-uploads" in storage.created_buckets
    assert any(u[0] == "b" * 64 + ".pdf" for u in bucket.uploads)


async def test_upload_pdf_to_storage_swallows_unexpected_error(monkeypatch):
    class _AngryBucket:
        def upload(self, *a, **k):
            raise RuntimeError("some transient 500")

    storage = _FakeStorage(_AngryBucket())
    monkeypatch.setattr(upload_service, "get_client", lambda use_admin=False: _FakeSB(storage))

    from backend.core import database

    async def _run_sync(fn, *a, **k):
        return fn(*a, **k)

    monkeypatch.setattr(database, "run_sync", _run_sync)

    # Must not raise — the worker retries storage; a failed upload here is logged.
    await upload_service.upload_pdf_to_storage("c" * 64, b"%PDF data")


# ── process_pdf_stream: parser routing labels (sync fallback vehicle) ─────────

async def _drain(agen, limit=50):
    chunks = []
    async for c in agen:
        chunks.append(c)
        if len(chunks) >= limit:
            break
    return chunks


@pytest.fixture
def sync_fallback(monkeypatch):
    """Force the dev in-process fallback and capture orchestrator kwargs."""
    from backend.core.config import settings
    from backend.services.parser import unified_orchestrator

    calls: dict = {}

    async def _no_storage(*_a, **_k):
        return None

    async def _arq_down(*_a, **_k):
        raise RuntimeError("redis down")

    async def _fake_unified(*_a, emit_fn=None, **kwargs):
        calls["unified"] = kwargs
        await emit_fn("complete", {"total": 1})

    monkeypatch.setattr(settings, "env", "development", raising=False)
    monkeypatch.setattr(settings, "parser_version", "5", raising=False)
    monkeypatch.setattr(upload_service, "upload_pdf_to_storage", _no_storage)
    monkeypatch.setattr(upload_service, "get_arq_pool", _arq_down)
    monkeypatch.setattr(unified_orchestrator, "parse_pdf_unified", _fake_unified)
    return calls


async def test_markitdown_branch_extracts_pages_and_converts(sync_fallback, monkeypatch):
    from backend.services import markitdown_service, office_convert

    async def _extract(content, filename):
        return {1: {"text": "slide one"}}

    async def _to_pdf(content, filename):
        return b"%PDF-converted"

    monkeypatch.setattr(markitdown_service, "extract_pages", _extract)
    monkeypatch.setattr(office_convert, "to_pdf", _to_pdf)

    await _drain(upload_service.process_pdf_stream(
        content=b"PK\x03\x04deck", filename="deck.pptx", pdf_hash="h" * 64,
        page_count=1, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser="auto", lecture_id=None, user_id="u1",
    ))
    # .pptx forces the markitdown parser regardless of requested parser.
    assert sync_fallback["unified"]["parser_used"] == "markitdown"
    assert sync_fallback["unified"]["odl_pages"] == {1: {"text": "slide one"}}


@pytest.mark.parametrize(
    "parser,module_name,expected_label",
    [
        ("llamaparse", "llamaparse_service", "llamaparse"),
        ("mineru", "mineru_service", "mineru"),
    ],
)
async def test_alternate_parser_branches(sync_fallback, monkeypatch, parser, module_name, expected_label):
    import importlib

    mod = importlib.import_module(f"backend.services.{module_name}")

    async def _extract(content, filename):
        return {1: {"text": f"from {parser}"}}

    monkeypatch.setattr(mod, "extract_pages", _extract)

    await _drain(upload_service.process_pdf_stream(
        content=b"%PDF data", filename="lecture.pdf", pdf_hash="h" * 64,
        page_count=1, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser=parser, lecture_id=None, user_id="u1",
    ))
    assert sync_fallback["unified"]["parser_used"] == expected_label
    assert sync_fallback["unified"]["odl_pages"] == {1: {"text": f"from {parser}"}}


async def test_opendataloader_branch(sync_fallback, monkeypatch):
    from backend.services import odl_service

    async def _extract(content, filename):
        return {1: {"text": "odl page"}}

    monkeypatch.setattr(odl_service, "extract_pages", _extract)

    await _drain(upload_service.process_pdf_stream(
        content=b"%PDF data", filename="lecture.pdf", pdf_hash="h" * 64,
        page_count=1, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser="opendataloader", lecture_id=None, user_id="u1",
    ))
    assert sync_fallback["unified"]["parser_used"] == "opendataloader-pdf"


async def test_course_and_visibility_threaded_to_unified(sync_fallback):
    await _drain(upload_service.process_pdf_stream(
        content=b"%PDF data", filename="lecture.pdf", pdf_hash="h" * 64,
        page_count=1, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser="pymupdf", lecture_id=None, user_id="u1",
        course_id="course-99", visibility="private",
    ))
    assert sync_fallback["unified"]["course_id"] == "course-99"
    assert sync_fallback["unified"]["visibility"] == "private"


# ── process_pdf_stream: production queue-outage refusal ───────────────────────

async def test_prod_queue_outage_yields_error_and_does_not_parse(monkeypatch):
    from backend.core.config import settings
    from backend.services.parser import unified_orchestrator

    ran = {"unified": False}

    async def _no_storage(*_a, **_k):
        return None

    async def _arq_down(*_a, **_k):
        raise RuntimeError("redis down")

    async def _fake_unified(*_a, **_k):
        ran["unified"] = True

    monkeypatch.setattr(settings, "env", "production", raising=False)
    monkeypatch.setattr(settings, "parser_version", "5", raising=False)
    monkeypatch.setattr(upload_service, "upload_pdf_to_storage", _no_storage)
    monkeypatch.setattr(upload_service, "get_arq_pool", _arq_down)
    monkeypatch.setattr(unified_orchestrator, "parse_pdf_unified", _fake_unified)

    chunks = await _drain(upload_service.process_pdf_stream(
        content=b"%PDF data", filename="lecture.pdf", pdf_hash="h" * 64,
        page_count=1, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser="pymupdf", lecture_id=None, user_id="u1",
    ))
    assert ran["unified"] is False
    assert any("temporarily unavailable" in c for c in chunks)
    assert any('"type": "error"' in c for c in chunks)


# ── process_pdf_stream: Arq enqueue + Redis pub/sub streaming ─────────────────

class _FakePubSub:
    def __init__(self, messages):
        self._messages = list(messages)

    async def __aenter__(self):
        return self

    async def __aexit__(self, *a):
        return False

    async def subscribe(self, channel):
        self.channel = channel

    async def get_message(self, ignore_subscribe_messages=False, timeout=None):
        if self._messages:
            return self._messages.pop(0)
        return None


class _FakeRedis:
    def __init__(self, messages):
        self._messages = messages
        self.closed = False

    def pubsub(self):
        return _FakePubSub(self._messages)

    async def aclose(self):
        self.closed = True


async def test_arq_path_enqueues_job_and_streams_events(monkeypatch):
    from backend.core.config import settings
    import redis.asyncio as aioredis

    enqueued = {}

    class _Pool:
        async def enqueue_job(self, name, **kwargs):
            enqueued["name"] = name
            enqueued["kwargs"] = kwargs

    async def _get_pool():
        return _Pool()

    async def _no_storage(*_a, **_k):
        return None

    messages = [
        None,  # heartbeat → ": ping"
        {"type": "subscribe"},  # ignored (not a "message")
        {"type": "message", "data": json.dumps({"type": "progress", "pct": 50})},
        {"type": "message", "data": json.dumps({"type": "complete", "total": 3})},
    ]
    fake_redis = _FakeRedis(messages)

    monkeypatch.setattr(settings, "parser_version", "5", raising=False)
    monkeypatch.setattr(upload_service, "get_arq_pool", _get_pool)
    monkeypatch.setattr(upload_service, "upload_pdf_to_storage", _no_storage)
    monkeypatch.setattr(aioredis, "from_url", lambda *a, **k: fake_redis)

    chunks = await _drain(upload_service.process_pdf_stream(
        content=b"%PDF data", filename="lecture.pdf", pdf_hash="d" * 64,
        page_count=3, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser="pymupdf", lecture_id=None, user_id="prof-7",
    ))

    assert enqueued["name"] == "parse_pdf_unified"
    assert enqueued["kwargs"]["pdf_hash"] == "d" * 64
    assert enqueued["kwargs"]["user_id"] == "prof-7"
    assert enqueued["kwargs"]["lecture_id"] == ""  # unified owns the lecture
    assert any(c == ": ping\n\n" for c in chunks)          # heartbeat emitted
    assert any('"progress"' in c for c in chunks)
    assert any('"complete"' in c for c in chunks)
    assert fake_redis.closed is True                        # redis client closed


async def test_arq_path_stops_on_error_event(monkeypatch):
    from backend.core.config import settings
    import redis.asyncio as aioredis

    class _Pool:
        async def enqueue_job(self, name, **kwargs):
            return None

    async def _get_pool():
        return _Pool()

    async def _no_storage(*_a, **_k):
        return None

    messages = [
        {"type": "message", "data": json.dumps({"type": "error", "message": "boom"})},
        {"type": "message", "data": json.dumps({"type": "complete"})},  # must NOT be reached
    ]
    fake_redis = _FakeRedis(messages)

    monkeypatch.setattr(settings, "parser_version", "5", raising=False)
    monkeypatch.setattr(upload_service, "get_arq_pool", _get_pool)
    monkeypatch.setattr(upload_service, "upload_pdf_to_storage", _no_storage)
    monkeypatch.setattr(aioredis, "from_url", lambda *a, **k: fake_redis)

    chunks = await _drain(upload_service.process_pdf_stream(
        content=b"%PDF data", filename="lecture.pdf", pdf_hash="e" * 64,
        page_count=1, ai_model="cerebras", use_blueprint=False,
        parsing_mode="ai", parser="pymupdf", lecture_id=None, user_id="u1",
    ))
    assert any('"error"' in c for c in chunks)
    assert not any('"complete"' in c for c in chunks)


# ── extract_raw_pages ────────────────────────────────────────────────────────

async def test_extract_raw_pages_pymupdf(sample_pdf_bytes):
    result = await upload_service.extract_raw_pages(sample_pdf_bytes, "lecture.pdf", "pymupdf")
    assert result["parser_used"] == "pymupdf"
    assert result["total_pages"] == 3
    first = result["pages"][0]
    assert first["page_num"] == 1
    assert "Slide 1 content" in first["text"]
    assert first["word_count"] == len(first["text"].split())
    assert first["char_count"] == len(first["text"])


@pytest.mark.parametrize(
    "parser,module_name,expected_label",
    [
        ("llamaparse", "llamaparse_service", "llamaparse"),
        ("mineru", "mineru_service", "mineru"),
        ("opendataloader", "odl_service", "opendataloader-pdf"),
    ],
)
async def test_extract_raw_pages_alternate_parsers(monkeypatch, parser, module_name, expected_label):
    import importlib

    mod = importlib.import_module(f"backend.services.{module_name}")

    async def _extract(content, filename):
        return {2: {"text": "second", "title": "T2"}, 1: {"text": "first", "title": "T1"}}

    monkeypatch.setattr(mod, "extract_pages", _extract)

    result = await upload_service.extract_raw_pages(b"%PDF data", "lecture.pdf", parser)
    assert result["parser_used"] == expected_label
    # Pages sorted ascending regardless of dict insertion order.
    assert [p["page_num"] for p in result["pages"]] == [1, 2]
    assert result["pages"][0]["title"] == "T1"


async def test_extract_raw_pages_auto_prefers_odl(monkeypatch):
    from backend.services import odl_service

    async def _extract(content, filename):
        return {1: {"text": "odl auto"}}

    monkeypatch.setattr(odl_service, "extract_pages", _extract)
    result = await upload_service.extract_raw_pages(b"%PDF data", "lecture.pdf", "auto")
    assert result["parser_used"] == "opendataloader-pdf"


async def test_extract_raw_pages_auto_falls_back_to_pymupdf(monkeypatch, sample_pdf_bytes):
    from backend.services import odl_service

    async def _extract(content, filename):
        raise RuntimeError("odl unavailable")

    monkeypatch.setattr(odl_service, "extract_pages", _extract)
    result = await upload_service.extract_raw_pages(sample_pdf_bytes, "lecture.pdf", "auto")
    assert result["parser_used"] == "pymupdf"
    assert result["total_pages"] == 3


# ── process_pdf_lazy ─────────────────────────────────────────────────────────

async def test_process_pdf_lazy_streams_updates(monkeypatch):
    async def _fake_lazy(content, filename=None, ai_model=None):
        yield {"type": "progress", "page": 1}
        yield {"type": "complete"}

    monkeypatch.setattr(upload_service, "import_pdf_lazy", _fake_lazy)

    chunks = await _drain(upload_service.process_pdf_lazy(b"%PDF data", "lecture.pdf", "cerebras"))
    assert any("pymupdf-lazy" in c for c in chunks)  # the info banner
    assert any('"complete"' in c for c in chunks)


async def test_process_pdf_lazy_emits_error_on_failure(monkeypatch):
    async def _fake_lazy(content, filename=None, ai_model=None):
        raise RuntimeError("lazy exploded")
        yield  # pragma: no cover — make it an async generator

    monkeypatch.setattr(upload_service, "import_pdf_lazy", _fake_lazy)

    chunks = await _drain(upload_service.process_pdf_lazy(b"%PDF data", "lecture.pdf", "cerebras"))
    assert any('"error"' in c and "lazy exploded" in c for c in chunks)
    assert any('"recoverable": false' in c for c in chunks)
